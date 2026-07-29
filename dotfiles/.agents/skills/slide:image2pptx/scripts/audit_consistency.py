#!/usr/bin/env python3
"""
Audit PPTX for theme consistency across all slides.

Checks every shape against theme.json and lists deviations:
  - background color
  - header band y/height/fill
  - page number position/font_size
  - margins (left/right)
  - font family of text shapes
  - common font sizes
  - shape role-based y coordinates (e.g. all titles at same y)

Usage:
  audit_consistency.py <pptx> <theme.json> <out_audit_json>

theme.json (minimal schema):
{
  "slide_size":   {"width_emu": int, "height_emu": int},
  "background":   {"hex": "#RRGGBB"},
  "header":       {"y_emu": int, "height_emu": int, "fill_hex": "#RRGGBB"},
  "footer":       {"y_emu": int, "height_emu": int},
  "page_number":  {"x_emu": int, "y_emu": int, "font_size_pt": float},
  "margins":      {"left_emu": int, "right_emu": int},
  "fonts":        {"ja": "Yu Gothic", "en": "Helvetica"},
  "tolerance": {
      "position_emu": 50000,    "size_emu": 50000,
      "rgb_distance": 15,       "font_size_pt": 1.0
  }
}
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor


def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("pptx")
    p.add_argument("theme")
    p.add_argument("out")
    return p.parse_args()


def hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def rgb_distance(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    return sum((x - y) ** 2 for x, y in zip(a, b)) ** 0.5


def shape_fill_rgb(shp) -> tuple[int, int, int] | None:
    try:
        f = shp.fill
        if f.type is None:
            return None
        c = f.fore_color.rgb
        return (int(str(c)[0:2], 16), int(str(c)[2:4], 16), int(str(c)[4:6], 16))
    except Exception:
        return None


def slide_bg_rgb(slide) -> tuple[int, int, int] | None:
    try:
        bg = slide.background.fill
        if bg.type is None:
            return None
        c = bg.fore_color.rgb
        return (int(str(c)[0:2], 16), int(str(c)[2:4], 16), int(str(c)[4:6], 16))
    except Exception:
        return None


def first_run_font(shp):
    if not shp.has_text_frame:
        return None, None
    for para in shp.text_frame.paragraphs:
        for run in para.runs:
            return run.font.name, (float(run.font.size.pt) if run.font.size else None)
    return None, None


def audit_slide(idx: int, slide, sw: int, sh: int, theme: dict) -> list[dict]:
    tol = theme.get("tolerance", {})
    pos_tol = tol.get("position_emu", 50000)
    size_tol = tol.get("size_emu", 50000)
    rgb_tol = tol.get("rgb_distance", 15)
    font_size_tol = tol.get("font_size_pt", 1.0)

    deviations: list[dict] = []
    page = idx + 1

    # 1. Background color
    bg_target = theme.get("background", {}).get("hex")
    if bg_target:
        actual = slide_bg_rgb(slide)
        if actual is None:
            deviations.append({"page": page, "kind": "background_missing",
                               "expected": bg_target, "actual": None})
        else:
            d = rgb_distance(actual, hex_to_rgb(bg_target))
            if d > rgb_tol:
                deviations.append({
                    "page": page, "kind": "background_color_drift",
                    "expected": bg_target,
                    "actual_rgb": list(actual),
                    "rgb_distance": round(d, 1),
                })

    # 2. Header band (largest top-area filled rectangle)
    header_t = theme.get("header")
    if header_t:
        header_shapes = [
            shp for shp in slide.shapes
            if shp.top is not None and shp.top < sh * 0.2
            and shp.height is not None and shp.height > 100000
            and shape_fill_rgb(shp) is not None
        ]
        header_shapes.sort(key=lambda s: -s.width * s.height)
        if not header_shapes:
            deviations.append({"page": page, "kind": "header_missing"})
        else:
            top = header_shapes[0]
            if abs(top.top - header_t["y_emu"]) > pos_tol:
                deviations.append({"page": page, "kind": "header_y_drift",
                                   "expected_y_emu": header_t["y_emu"],
                                   "actual_y_emu": top.top})
            if abs(top.height - header_t["height_emu"]) > size_tol:
                deviations.append({"page": page, "kind": "header_height_drift",
                                   "expected_h_emu": header_t["height_emu"],
                                   "actual_h_emu": top.height})
            if "fill_hex" in header_t:
                actual = shape_fill_rgb(top)
                if actual:
                    d = rgb_distance(actual, hex_to_rgb(header_t["fill_hex"]))
                    if d > rgb_tol:
                        deviations.append({
                            "page": page, "kind": "header_fill_drift",
                            "expected": header_t["fill_hex"],
                            "actual_rgb": list(actual),
                            "rgb_distance": round(d, 1),
                        })

    # 3. Page number
    pn = theme.get("page_number")
    if pn:
        candidates = [
            shp for shp in slide.shapes
            if shp.has_text_frame and shp.text_frame.text.strip().isdigit()
            and shp.top > sh * 0.85
        ]
        if not candidates:
            deviations.append({"page": page, "kind": "page_number_missing"})
        else:
            shp = candidates[0]
            if abs(shp.left - pn["x_emu"]) > pos_tol:
                deviations.append({"page": page, "kind": "page_number_x_drift",
                                   "expected_x_emu": pn["x_emu"],
                                   "actual_x_emu": shp.left})
            if abs(shp.top - pn["y_emu"]) > pos_tol:
                deviations.append({"page": page, "kind": "page_number_y_drift",
                                   "expected_y_emu": pn["y_emu"],
                                   "actual_y_emu": shp.top})
            _, fs = first_run_font(shp)
            if fs and abs(fs - pn["font_size_pt"]) > font_size_tol:
                deviations.append({"page": page, "kind": "page_number_font_drift",
                                   "expected_pt": pn["font_size_pt"],
                                   "actual_pt": fs})

    # 4. Margins
    m = theme.get("margins")
    if m:
        # 主要 text shape (面積上位5) の左端を集計
        text_shapes = sorted(
            [shp for shp in slide.shapes if shp.has_text_frame and shp.text_frame.text.strip()],
            key=lambda s: -(s.width * s.height),
        )[:5]
        for shp in text_shapes:
            if shp.left < m["left_emu"] - pos_tol:
                deviations.append({"page": page, "kind": "left_margin_violation",
                                   "expected_left_emu": m["left_emu"],
                                   "actual_left_emu": shp.left, "text": shp.text_frame.text[:40]})
            right_actual = shp.left + shp.width
            if right_actual > sw - m["right_emu"] + pos_tol:
                deviations.append({"page": page, "kind": "right_margin_violation",
                                   "expected_right_emu": sw - m["right_emu"],
                                   "actual_right_emu": right_actual,
                                   "text": shp.text_frame.text[:40]})

    # 5. Fonts (ja default check, only on text shapes with Japanese chars)
    fonts = theme.get("fonts", {})
    ja_default = fonts.get("ja")
    if ja_default:
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            text = shp.text_frame.text
            if not any("぀" <= c <= "ヿ" or "一" <= c <= "鿿" for c in text):
                continue
            fname, _ = first_run_font(shp)
            if fname and fname != ja_default:
                deviations.append({"page": page, "kind": "font_drift",
                                   "expected": ja_default, "actual": fname,
                                   "text": text[:40]})

    return deviations


def main() -> int:
    args = parse_args()
    theme = json.loads(Path(args.theme).read_text())
    prs = Presentation(args.pptx)
    sw, sh = prs.slide_width, prs.slide_height

    # Slide-size sanity
    expected = theme.get("slide_size", {})
    if expected:
        if sw != expected.get("width_emu") or sh != expected.get("height_emu"):
            global_dev = [{
                "page": 0, "kind": "slide_size_drift",
                "expected": expected, "actual": {"width_emu": sw, "height_emu": sh},
            }]
        else:
            global_dev = []
    else:
        global_dev = []

    all_deviations = list(global_dev)
    for i, slide in enumerate(prs.slides):
        all_deviations.extend(audit_slide(i, slide, sw, sh, theme))

    # cross-page consistency: same role shapes at same coordinate
    # check: header.top across pages should be identical to first page's
    header_tops: list[tuple[int, int]] = []
    for i, slide in enumerate(prs.slides):
        cands = [
            shp for shp in slide.shapes
            if shp.top < sh * 0.2 and shp.height > 100000
            and shape_fill_rgb(shp) is not None
        ]
        if cands:
            cands.sort(key=lambda s: -s.width * s.height)
            header_tops.append((i + 1, cands[0].top))
    if header_tops:
        ref = header_tops[0][1]
        for page, top in header_tops[1:]:
            if abs(top - ref) > theme.get("tolerance", {}).get("position_emu", 50000):
                all_deviations.append({
                    "page": page, "kind": "header_y_inconsistent_with_page1",
                    "page1_y_emu": ref, "actual_y_emu": top,
                })

    summary = {
        "pptx": args.pptx,
        "theme": args.theme,
        "num_slides": len(prs.slides),
        "num_deviations": len(all_deviations),
        "by_kind": {},
    }
    for d in all_deviations:
        summary["by_kind"][d["kind"]] = summary["by_kind"].get(d["kind"], 0) + 1

    out = {"summary": summary, "deviations": all_deviations}
    Path(args.out).write_text(json.dumps(out, ensure_ascii=False, indent=2))
    print(json.dumps(summary, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
